from flask import Flask, render_template, request, redirect, url_for, send_from_directory
import openai
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
from dotenv import load_dotenv
import os

load_dotenv()

openai.api_key = os.getenv('API_KEY')

TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

app = Flask(__name__)

def create_slide_titles(topic, num_slides, presentation_type, extra_details):
    prompt = f"Generate {num_slides} short slide titles for a '{presentation_type}' presentation on the topic '{topic}'. {extra_details}"
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=200,
        n=1
    )
    return response['choices'][0]['message']['content'].split("\n")

def create_slide_content(slide_title, presentation_type, extra_details):
    prompt = f"Generate content for the slide: '{slide_title}'. The content must be in medium-worded paragraphs. Include details for a '{presentation_type}' presentation. {extra_details}. Only return 2 paragraphs."
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=300,
        n=1
    )
    return response['choices'][0]['message']['content']

def create_presentation(topic, slide_titles, slide_contents, creator_name, theme):
    powerpoint = Presentation()

    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    themes = {
        "light": {"bg": RGBColor(255, 255, 255), "font": RGBColor(0, 0, 0)},
        "dark": {"bg": RGBColor(0, 0, 0), "font": RGBColor(255, 255, 255)},
        "blue": {"bg": RGBColor(0, 0, 255), "font": RGBColor(255, 255, 255)}
    }

    selected_theme = themes.get(theme, {"bg": RGBColor(173, 216, 230), "font": RGBColor(0, 0, 0)})

    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic

    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = selected_theme["font"]
    content = title_slide.placeholders[1]
    content.text = "Created By " + creator_name
    content.text_frame.paragraphs[0].font.size = Pt(24)
    content.text_frame.paragraphs[0].font.color.rgb = selected_theme["font"]

    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = selected_theme["bg"]

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = powerpoint.slides.add_slide(content_slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = selected_theme["bg"]

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = selected_theme["font"]

        content = slide.placeholders[1]
        content.text = slide_content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = SLIDE_FONT_SIZE
            paragraph.font.color.rgb = selected_theme["font"]

    if not os.path.exists('powerpoints'):
        os.makedirs('powerpoints')

    file_path = f"powerpoints/{topic}.pptx"
    powerpoint.save(file_path)
    return file_path

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        topic = request.form['topic']
        num_slides = int(request.form['num_slides'])
        presentation_type = request.form['presentation_type']
        extra_details = request.form['extra_details']
        creator_name = request.form['creator_name']
        theme = request.form['theme']

        slide_titles = create_slide_titles(topic, num_slides, presentation_type, extra_details)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        slide_contents = [create_slide_content(title, presentation_type, extra_details) for title in filtered_slide_titles]

        file_path = create_presentation(topic, filtered_slide_titles, slide_contents, creator_name, theme)

        return redirect(url_for('download_file', filename=os.path.basename(file_path)))

    return render_template('index.html')

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory('powerpoints', filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
